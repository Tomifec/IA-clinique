#!/usr/bin/env python3
"""
Extract DOIs from heterogeneous local sources (offline):

- Supported inputs: directories or files (.md, .txt, .html/.htm, .pdf, .pptx)
- No network access. Pure regex on text. For .pdf/.pptx, tries optional libs.
- Outputs:
  - artifacts/imports/extracted_dois_sources_<date>.txt (unique DOIs)
  - artifacts/imports/extracted_dois_sources_<date>.csv (file,doi pairs)

Usage (Windows PowerShell examples):
  python scripts/extract_dois_from_sources.py C:\path\to\folder more\files\notes.pdf
  # Then seed via agents:
  #   $dois = Get-Content artifacts\imports\extracted_dois_sources_YYYY-MM-DD.txt
  #   python tools/agents_online/seed_by_doi.py --topic import_refs_YYYY_MM_DD $dois

Notes:
- .pdf: uses PyPDF2 if installed; otherwise skips with a warning.
- .pptx: uses python-pptx if installed; otherwise skips with a warning.
"""
from __future__ import annotations

import argparse
import datetime as dt
import pathlib
import re
import sys
from typing import Iterable, List, Set, Tuple


DOI_RX = re.compile(r"10\.\d{4,9}/\S+", re.I)
DOI_URL_RX = re.compile(r"https?://(?:dx\.)?doi\.org/(10\.\d{4,9}/\S+)", re.I)


def extract_from_text(text: str) -> Set[str]:
    out: Set[str] = set()
    for m in DOI_URL_RX.finditer(text):
        out.add(m.group(1))
    for m in DOI_RX.finditer(text):
        out.add(m.group(0))
    cleaned: Set[str] = set()
    for d in out:
        # Trim common enclosing characters and trailing punctuation/quotes
        x = d.strip().strip("<>[](){}\"' ")
        x = x.rstrip(").,;:*\u201d\u2019\u00bb")
        x = x.replace(" ", "")
        if x:
            cleaned.add(x)
    return cleaned


def read_text_file(p: pathlib.Path) -> str:
    try:
        return p.read_text(encoding="utf-8", errors="ignore")
    except Exception:
        return ""


def read_pdf_text(p: pathlib.Path, max_pages: int | None = 30) -> str:
    try:
        import PyPDF2  # type: ignore
    except Exception:
        sys.stderr.write(f"[WARN] PyPDF2 not installed, skip PDF: {p}\n")
        return ""
    try:
        txt_parts: List[str] = []
        with open(p, "rb") as f:
            reader = PyPDF2.PdfReader(f)
            n = len(reader.pages)
            limit = n if max_pages is None else min(n, max_pages)
            for i in range(limit):
                try:
                    page = reader.pages[i]
                    t = page.extract_text() or ""
                except Exception:
                    t = ""
                if t:
                    txt_parts.append(t)
        return "\n".join(txt_parts)
    except Exception:
        return ""


def read_pptx_text(p: pathlib.Path) -> str:
    try:
        from pptx import Presentation  # type: ignore
    except Exception:
        sys.stderr.write(f"[WARN] python-pptx not installed, skip PPTX: {p}\n")
        return ""
    try:
        prs = Presentation(str(p))
        parts: List[str] = []
        for s in prs.slides:
            for shape in s.shapes:
                try:
                    if hasattr(shape, "text") and shape.text:
                        parts.append(shape.text)
                except Exception:
                    continue
        return "\n".join(parts)
    except Exception:
        return ""


def iter_files(paths: Iterable[str]) -> Iterable[pathlib.Path]:
    for raw in paths:
        p = pathlib.Path(raw)
        if p.is_dir():
            for f in p.rglob("*"):
                if f.is_file():
                    yield f
        elif p.is_file():
            yield p


def main(argv: List[str] | None = None) -> int:
    ap = argparse.ArgumentParser(description="Extract DOIs from local sources (offline)")
    ap.add_argument("paths", nargs="+", help="Files or directories to scan")
    ap.add_argument("--max-pdf-pages", type=int, default=30, help="Max pages to read per PDF (default 30)")
    args = ap.parse_args(argv)

    pairs: List[Tuple[str, str]] = []
    all_dois: Set[str] = set()

    exts_text = {".md", ".txt", ".html", ".htm"}
    for f in iter_files(args.paths):
        ext = f.suffix.lower()
        text = ""
        if ext in exts_text:
            text = read_text_file(f)
        elif ext == ".pdf":
            text = read_pdf_text(f, max_pages=args.max_pdf_pages)
        elif ext == ".pptx":
            text = read_pptx_text(f)
        else:
            continue
        if not text:
            continue
        dois = extract_from_text(text)
        for d in sorted(dois):
            pairs.append((str(f), d))
        all_dois.update(dois)

    out_dir = pathlib.Path("artifacts") / "imports"
    out_dir.mkdir(parents=True, exist_ok=True)
    stamp = dt.date.today().isoformat()
    txt_path = out_dir / f"extracted_dois_sources_{stamp}.txt"
    csv_path = out_dir / f"extracted_dois_sources_{stamp}.csv"
    txt_path.write_text("\n".join(sorted(all_dois)) + "\n", encoding="utf-8")
    csv_lines = ["file,doi"] + [f"{pair[0]},{pair[1]}" for pair in pairs]
    csv_path.write_text("\n".join(csv_lines) + "\n", encoding="utf-8")

    print(f"files_scanned={len(set(p for p,_ in pairs))} unique_dois={len(all_dois)}")
    print(f"txt -> {txt_path}")
    print(f"csv -> {csv_path}")
    return 0


if __name__ == "__main__":
    raise SystemExit(main())

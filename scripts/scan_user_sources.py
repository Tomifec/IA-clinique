#!/usr/bin/env python3
import argparse
import datetime as dt
import json
import os
import pathlib
import re
import subprocess
import sys
from typing import Dict, List, Set, Tuple


# Reuse extractor
try:
    from scripts.extract_dois_from_sources import extract_from_text
except Exception:
    # When run as `python scripts/scan_user_sources.py`
    sys.path.append(str(pathlib.Path(__file__).resolve().parents[1]))
    from scripts.extract_dois_from_sources import extract_from_text  # type: ignore

try:
    from tools.agents_online.curation_utils import read_fm
except Exception:
    from curation_utils import read_fm  # type: ignore


TEXT_EXTS = {".md", ".txt", ".html", ".htm"}


def read_text(path: pathlib.Path) -> str:
    try:
        if path.suffix.lower() in TEXT_EXTS:
            return path.read_text(encoding="utf-8", errors="ignore")
        if path.suffix.lower() == ".pdf":
            try:
                import PyPDF2  # type: ignore
                txt_parts: List[str] = []
                with open(path, "rb") as f:
                    reader = PyPDF2.PdfReader(f)
                    n = len(reader.pages)
                    for i in range(min(30, n)):
                        try:
                            t = reader.pages[i].extract_text() or ""
                            if t:
                                txt_parts.append(t)
                        except Exception:
                            continue
                return "\n".join(txt_parts)
            except Exception:
                return ""
        if path.suffix.lower() == ".pptx":
            try:
                from pptx import Presentation  # type: ignore
                prs = Presentation(str(path))
                parts: List[str] = []
                for s in prs.slides:
                    for shape in s.shapes:
                        try:
                            if hasattr(shape, "text") and shape.text:
                                parts.append(shape.text)
                        except Exception:
                            continue
                return "\n".join(parts)
            except Exception:
                return ""
    except Exception:
        return ""
    return ""


def existing_dois(evidence_dir: pathlib.Path) -> Set[str]:
    out: Set[str] = set()
    for md in evidence_dir.rglob("*.md"):
        try:
            fm, _ = read_fm(md)
            du = fm.get("doi_url")
            if isinstance(du, str) and du.strip():
                m = re.search(r"(10\.\d{4,9}/\S+)", du)
                if m:
                    out.add(m.group(1))
        except Exception:
            continue
    return out


def main(argv: List[str] | None = None) -> int:
    ap = argparse.ArgumentParser(description="Scan progressif des sources utilisateur et extraction de DOIs")
    ap.add_argument("--root", default=str(pathlib.Path.home()), help="Racine à scanner (par défaut: HOME)")
    ap.add_argument("--limit", type=int, default=50, help="Nombre max de NOUVEAUX DOIs à collecter ce tour")
    ap.add_argument("--state", default="logs/sources_scan_state.json", help="Fichier d'état (JSON)")
    ap.add_argument("--seed", action="store_true", help="Seed les Evidence depuis les nouveaux DOIs")
    ap.add_argument("--curate", action="store_true", help="Lance la curation après seed")
    ap.add_argument("--topic", default=f"import_refs_{dt.date.today().isoformat().replace('-', '_')}")
    args = ap.parse_args(argv)

    root = pathlib.Path(args.root)
    state_path = pathlib.Path(args.state)
    state: Dict[str, Dict[str, float]] = {"processed": {}}
    if state_path.exists():
        try:
            state = json.loads(state_path.read_text(encoding="utf-8"))
        except Exception:
            state = {"processed": {}}

    processed: Dict[str, float] = state.get("processed", {})
    already: Set[str] = existing_dois(pathlib.Path("evidence"))

    new_pairs: List[Tuple[str, str]] = []
    new_dois: List[str] = []

    exts = TEXT_EXTS | {".pdf", ".pptx"}
    for p in root.rglob("*"):
        if not p.is_file():
            continue
        if p.suffix.lower() not in exts:
            continue
        try:
            mtime = p.stat().st_mtime
        except Exception:
            continue
        key = str(p)
        if processed.get(key) and processed[key] >= mtime:
            continue  # already scanned up-to-date
        text = read_text(p)
        if not text:
            processed[key] = mtime
            continue
        dois = sorted(extract_from_text(text))
        added_here = False
        for d in dois:
            if d in already or d in new_dois:
                continue
            new_pairs.append((key, d))
            new_dois.append(d)
            added_here = True
            if len(new_dois) >= args.limit:
                break
        processed[key] = mtime
        if len(new_dois) >= args.limit:
            break

    out_dir = pathlib.Path("artifacts") / "imports"
    out_dir.mkdir(parents=True, exist_ok=True)
    stamp = dt.date.today().isoformat()
    txt_path = out_dir / f"extracted_dois_sources_{stamp}.txt"
    csv_path = out_dir / f"extracted_dois_sources_{stamp}.csv"
    if new_dois:
        with txt_path.open("a", encoding="utf-8") as f:
            for d in new_dois:
                f.write(d + "\n")
        with csv_path.open("a", encoding="utf-8") as f:
            if csv_path.stat().st_size == 0 if csv_path.exists() else True:
                f.write("file,doi\n")
            for file, doi in new_pairs:
                f.write(f"{file},{doi}\n")

    # Persist state
    state["processed"] = processed
    state["last_run"] = dt.datetime.now().isoformat()
    state["last_batch"] = {"count": len(new_dois)}
    state_path.parent.mkdir(parents=True, exist_ok=True)
    state_path.write_text(json.dumps(state, ensure_ascii=False, indent=2), encoding="utf-8")

    print(f"scanned_new_files={len(new_pairs)} new_dois={len(new_dois)}")
    if new_dois:
        print(f"txt -> {txt_path}")
        print(f"csv -> {csv_path}")

    # Optional seed + curate
    if new_dois and args.seed:
        cmd = [
            sys.executable,
            "tools/agents_online/seed_by_doi.py",
            "--topic",
            args.topic,
            *new_dois,
        ]
        try:
            print("Running:", " ".join(cmd))
        except UnicodeEncodeError:
            safe = " ".join(cmd)
            try:
                safe = safe.encode("cp1252", "replace").decode("cp1252")
            except Exception:
                safe = safe.encode("ascii", "replace").decode("ascii")
            print("Running:", safe)
        subprocess.run(cmd, check=False)
        if args.curate:
            subprocess.run([sys.executable, "tools/agents_online/curate.py", "--write", "evidence"], check=False)

    return 0


if __name__ == "__main__":
    raise SystemExit(main())
